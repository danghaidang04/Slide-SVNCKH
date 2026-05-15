"""Tìm các ký tự tam giác đỏ (▶ / blacktriangleright) trong PPTX gốc."""
from pptx import Presentation

prs = Presentation("Slide_HUST.pptx")

NSMAP_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

# Các ký tự Unicode có thể là tam giác
TRIANGLE_CHARS = {'▶', '▷', '►', '▸', '▹', '◀', '◁', '◂', '◃', '◄', 
                  '\u25B6', '\u25B7', '\u25BA', '\u25B8', '\u25B9',
                  '\u25C0', '\u25C1', '\u25C2', '\u25C3', '\u25C4',
                  'I'}  # Adobe đôi khi render thành "I"

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run in para.runs:
                text = run.text.strip()
                if not text:
                    continue
                
                font_name = run.font.name or "(inherited)"
                font_size = f"{run.font.size.pt:.1f}pt" if run.font.size else "(inh)"
                
                try:
                    color = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else "(inh)"
                except:
                    color = "(inh)"
                
                # Kiểm tra các ký tự đặc biệt
                has_special = False
                for ch in text:
                    code = ord(ch)
                    if ch in TRIANGLE_CHARS or (0x2500 <= code <= 0x27FF) or (0xE000 <= code <= 0xF8FF):
                        has_special = True
                        break
                
                if has_special or len(text) == 1:
                    hex_chars = ' '.join(f'U+{ord(c):04X}' for c in text)
                    print(f"Slide {slide_idx+1} | {shape.name} | Font: {font_name} | Size: {font_size} | Color: {color}")
                    print(f"  Text: '{text}' | Unicode: {hex_chars}")
                    print()
