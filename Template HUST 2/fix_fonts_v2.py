from pptx import Presentation
from pptx.util import Pt

# Thử các biến thể tên font phổ biến trên Mac
FONT_SANS = "LM Sans 10"
FONT_MONO = "LM Mono 10"

def apply_font_to_shape(shape, is_header):
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_SANS
            
            # Thiết lập size dựa trên vị trí (Header vs Body)
            if is_header:
                run.font.size = Pt(16)
                run.font.bold = True
            else:
                run.font.size = Pt(12)

def process_recursive(shapes, slide_height):
    for shape in shapes:
        if shape.shape_type == 6: # Group
            process_recursive(shape.shapes, slide_height)
        elif shape.has_text_frame:
            # Nếu nằm ở 20% phía trên slide thì coi là Header
            is_header = shape.top < (slide_height * 0.2)
            apply_font_to_shape(shape, is_header)

def fix_fonts_v2(input_path, output_path):
    prs = Presentation(input_path)
    slide_height = prs.slide_height
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Fixing Slide {slide_idx + 1}...")
        
        # Trang đầu tiên xử lý đặc biệt (Title Slide)
        if slide_idx == 0:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = FONT_SANS
                            # Tiêu đề cực lớn
                            if shape.top < (slide_height * 0.5):
                                run.font.size = Pt(20)
                                run.font.bold = True
                            else:
                                run.font.size = Pt(11)
        else:
            process_recursive(slide.shapes, slide_height)
            
    prs.save(output_path)
    print(f"Hoàn thành! Đã lưu bản v2 tại {output_path}")

if __name__ == "__main__":
    fix_fonts_v2("Slide_HUST.pptx", "Slide_HUST_Fixed_v2.pptx")
