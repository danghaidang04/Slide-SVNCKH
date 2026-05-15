from pptx import Presentation
from pptx.util import Pt

# Font chuẩn Mac
FONT_SANS = "LM Sans 10"

# Hệ số tỷ lệ từ Beamer sang PowerPoint (Xấp xỉ 1.6)
SCALE = 1.6

def fix_fonts_v4(input_path, output_path):
    prs = Presentation(input_path)
    slide_height = prs.slide_height
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Scaling and Fixing Slide {slide_idx + 1}...")
        
        for shape in slide.shapes:
            if shape.shape_type == 6: # Group
                shapes_to_process = shape.shapes
            else:
                shapes_to_process = [shape]
                
            for s in shapes_to_process:
                if not s.has_text_frame:
                    continue
                
                rel_top = s.top / slide_height
                
                for paragraph in s.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONT_SANS
                        
                        if slide_idx == 0:
                            # Title Slide (Hệ số 1.6)
                            if rel_top < 0.5: # Main Title (18pt * 1.6 = 28.8)
                                run.font.size = Pt(28)
                                run.font.bold = True
                            else: # Info (10pt * 1.6 = 16)
                                run.font.size = Pt(16)
                        else:
                            # Content Slides
                            if rel_top < 0.15: # Frame Title (16pt * 1.6 = 25.6)
                                run.font.size = Pt(25)
                                run.font.bold = True
                            else: # Body text (11pt * 1.6 = 17.6)
                                run.font.size = Pt(18)
                                
    prs.save(output_path)
    print(f"Đã xong bản v4 (Đã Scale cỡ chữ) tại {output_path}")

if __name__ == "__main__":
    fix_fonts_v4("Slide_HUST.pptx", "Slide_HUST_Fixed_v4.pptx")
