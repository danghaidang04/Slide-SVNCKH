"""
Fix TOÀN BỘ font trong PPTX — bao gồm cả:
- Runs thông thường (text)
- Runs inherited (không có font rõ ràng)
- Slide Master, Layout, Theme
- PLACEHOLDER số trang (field elements)
- Character spacing
- Tam giác đỏ ▶
- Chấm tròn xanh ●
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree

FONTS_TO_REPLACE = {"Arial Black", "Arial", "Calibri", "Times New Roman"}
TARGET_FONT = "LM Sans 10"
NSMAP_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NSMAP_A_BRACKET = f'{{{NSMAP_A}}}'

stats = {
    'font': 0, 'font_inherited': 0, 'spacing': 0,
    'triangle': 0, 'bullet': 0, 'master': 0, 'pagenum': 0
}

def set_latin_font_on_rPr(rPr_elem):
    """Thêm hoặc sửa <a:latin typeface="LM Sans 10"/> trong rPr."""
    latin = rPr_elem.find(f'{NSMAP_A_BRACKET}latin')
    if latin is not None:
        if latin.get('typeface', '') != TARGET_FONT:
            latin.set('typeface', TARGET_FONT)
            return True
    else:
        latin = etree.SubElement(rPr_elem, f'{NSMAP_A_BRACKET}latin')
        latin.set('typeface', TARGET_FONT)
        return True
    return False

def fix_runs_in_textframe(tf):
    """Fix tất cả runs + field elements trong text frame."""
    for paragraph in tf.paragraphs:
        # Fix paragraph-level bullets
        pPr = paragraph._pPr
        if pPr is not None:
            buFont = pPr.find(f'{NSMAP_A_BRACKET}buFont')
            buSzPct = pPr.find(f'{NSMAP_A_BRACKET}buSzPct')
            buChar = pPr.find(f'{NSMAP_A_BRACKET}buChar')
            
            if buChar is not None:
                if buFont is not None:
                    old = buFont.get('typeface', '')
                    if old != TARGET_FONT:
                        buFont.set('typeface', TARGET_FONT)
                
                if buSzPct is not None:
                    old_sz = int(buSzPct.get('val', '100000'))
                    if old_sz < 100000:
                        buSzPct.set('val', '100000')
                        stats['bullet'] += 1
        
        # Fix text runs
        for run in paragraph.runs:
            current_font = run.font.name
            
            if current_font in FONTS_TO_REPLACE:
                run.font.name = TARGET_FONT
                stats['font'] += 1
            elif current_font is None:
                run.font.name = TARGET_FONT
                stats['font_inherited'] += 1
            
            # Character spacing
            rPr = run._r.find(f'{NSMAP_A_BRACKET}rPr')
            if rPr is not None and 'spc' in rPr.attrib:
                spc_val = int(rPr.get('spc'))
                
                is_red_triangle = False
                if spc_val > 400 and run.text.strip() == "I":
                    try:
                        if run.font.color and run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00):
                            is_red_triangle = True
                    except:
                        pass
                
                if is_red_triangle:
                    run.text = "▶"
                    del rPr.attrib['spc']
                    stats['triangle'] += 1
                elif spc_val < 0:
                    del rPr.attrib['spc']
                    stats['spacing'] += 1
    
    # Fix FIELD elements (số trang, ngày tháng, etc.)
    p_elem = tf._txBody
    for fld in p_elem.iter(f'{NSMAP_A_BRACKET}fld'):
        rPr = fld.find(f'{NSMAP_A_BRACKET}rPr')
        if rPr is not None:
            if set_latin_font_on_rPr(rPr):
                stats['pagenum'] += 1
            # Xóa spacing âm trên số trang
            if 'spc' in rPr.attrib:
                spc_val = int(rPr.get('spc'))
                if spc_val < 0:
                    del rPr.attrib['spc']

def fix_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:  # Group
            fix_shapes(shape.shapes)
        elif shape.has_text_frame:
            fix_runs_in_textframe(shape.text_frame)

def fix_master_and_layouts(prs):
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name != TARGET_FONT:
                            run.font.name = TARGET_FONT
                            stats['master'] += 1
                    pPr = para._pPr
                    if pPr is not None:
                        for defRPr in pPr.findall(f'{NSMAP_A_BRACKET}defRPr'):
                            if set_latin_font_on_rPr(defRPr):
                                stats['master'] += 1
        
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name != TARGET_FONT:
                                run.font.name = TARGET_FONT
                                stats['master'] += 1
                        pPr = para._pPr
                        if pPr is not None:
                            for defRPr in pPr.findall(f'{NSMAP_A_BRACKET}defRPr'):
                                if set_latin_font_on_rPr(defRPr):
                                    stats['master'] += 1
            
            # Fix field elements in layouts too
            for fld in layout._element.iter(f'{NSMAP_A_BRACKET}fld'):
                rPr = fld.find(f'{NSMAP_A_BRACKET}rPr')
                if rPr is not None:
                    if set_latin_font_on_rPr(rPr):
                        stats['master'] += 1

def fix_theme_fonts(prs):
    for master in prs.slide_masters:
        for elem in master.element.iter():
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if tag in ('latin', 'ea', 'cs'):
                old = elem.get('typeface', '')
                if old and old != TARGET_FONT and old not in ('+mj-lt', '+mn-lt', '+mj-ea', '+mn-ea'):
                    elem.set('typeface', TARGET_FONT)
                    stats['master'] += 1

def main():
    prs = Presentation("Slide_HUST.pptx")
    
    print("Fixing Slide Master, Layouts, Theme...")
    fix_master_and_layouts(prs)
    fix_theme_fonts(prs)
    
    for slide_idx, slide in enumerate(prs.slides):
        fix_shapes(slide.shapes)
        print(f"  Slide {slide_idx + 1}: OK")
    
    output = "Slide_HUST_Fixed_Final.pptx"
    prs.save(output)
    
    print(f"\n{'='*50}")
    print(f"KẾT QUẢ TỔNG HỢP:")
    print(f"{'='*50}")
    print(f"  Font đổi (explicit):    {stats['font']} runs")
    print(f"  Font đổi (inherited):   {stats['font_inherited']} runs")
    print(f"  Số trang (field):       {stats['pagenum']} fields")
    print(f"  Spacing âm xóa:         {stats['spacing']} runs")
    print(f"  Tam giác đỏ ▶:          {stats['triangle']} runs")
    print(f"  Chấm tròn xanh ●:      {stats['bullet']} bullets")
    print(f"  Master/Layout/Theme:     {stats['master']} elements")
    print(f"{'='*50}")
    print(f"  → Lưu tại: {output}")

if __name__ == "__main__":
    main()
