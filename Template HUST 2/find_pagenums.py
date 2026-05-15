"""Kiểm tra XML của Placeholder số trang + tìm cách fix."""
from pptx import Presentation
from lxml import etree

NSMAP_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
TARGET_FONT = "LM Sans 10"

# Kiểm tra bản FIXED
prs = Presentation("Slide_HUST_Fixed_Final.pptx")

slide = prs.slides[1]  # Slide 2 (0-indexed)
for shape in slide.shapes:
    if shape.left and shape.left / 360000 > 14 and shape.top / 360000 > 7:
        print(f"Shape: '{shape.name}' | Type: {shape.shape_type}")
        print(f"has_text_frame: {shape.has_text_frame}")
        print(f"placeholder_format: {shape.placeholder_format}")
        if shape.placeholder_format:
            print(f"  idx: {shape.placeholder_format.idx}")
            print(f"  type: {shape.placeholder_format.type}")
        
        # In XML
        xml_str = etree.tostring(shape._element, pretty_print=True).decode()
        print(f"\nFull XML:\n{xml_str}")
        
        # Thử truy cập text frame bằng cách khác
        txBody = shape._element.find(f'{NSMAP_A}txBody')
        if txBody is None:
            txBody = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}txBody')
        if txBody is None:
            # Tìm trong namespace p
            for child in shape._element:
                print(f"  Child tag: {child.tag}")
        
        break
