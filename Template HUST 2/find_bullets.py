"""Tìm chấm tròn: kiểm tra paragraph bullet level và shape nhỏ."""
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

prs = Presentation("Slide_HUST.pptx")
NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# Slides có bullet: 5 (evaluation) và 7 (conclusions)
for slide_idx in [4, 6]:  # 0-indexed
    slide = prs.slides[slide_idx]
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_idx+1}")
    print(f"{'='*60}")
    
    for shape in slide.shapes:
        # Kiểm tra shapes rất nhỏ (có thể là chấm tròn)
        if not shape.has_text_frame:
            w = shape.width / 360000 if shape.width else 0
            h = shape.height / 360000 if shape.height else 0
            if w < 0.5 and h < 0.5 and w > 0:
                top = shape.top / 360000 if shape.top else 0
                left = shape.left / 360000 if shape.left else 0
                print(f"  [SMALL SHAPE] '{shape.name}' | Type: {shape.shape_type} | Pos: ({left:.2f}, {top:.2f}) | Size: {w:.2f}x{h:.2f} cm")
            continue
        
        # Kiểm tra paragraph-level bullet properties
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            pPr = para._pPr
            if pPr is not None:
                buNone = pPr.find(f'{{{NSMAP["a"]}}}buNone')
                buChar = pPr.find(f'{{{NSMAP["a"]}}}buChar')
                buFont = pPr.find(f'{{{NSMAP["a"]}}}buFont')
                buSzPct = pPr.find(f'{{{NSMAP["a"]}}}buSzPct')
                buClr = pPr.find(f'{{{NSMAP["a"]}}}buClr')
                lvl = pPr.get('lvl', '0')
                
                if buChar is not None:
                    char_val = buChar.get('char', '?')
                    font_val = buFont.get('typeface', '?') if buFont is not None else '?'
                    sz_val = buSzPct.get('val', '?') if buSzPct is not None else '?'
                    
                    # Get bullet color
                    clr_val = "?"
                    if buClr is not None:
                        srgb = buClr.find(f'{{{NSMAP["a"]}}}srgbClr')
                        if srgb is not None:
                            clr_val = srgb.get('val', '?')
                    
                    text = ''.join(r.text for r in para.runs)[:40]
                    print(f"  [BULLET] lvl={lvl} char='{char_val}' font={font_val} size={sz_val} color={clr_val}")
                    print(f"    Text: '{text}'")
