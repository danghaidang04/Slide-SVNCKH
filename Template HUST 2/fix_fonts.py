from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Tên font chuẩn đã cài trên Mac
FONT_SANS = "LM Sans 10"
FONT_ROMAN = "LM Roman 10"

def fix_fonts(input_path, output_path):
    prs = Presentation(input_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx + 1}...")
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            # Xác định đây là tiêu đề hay nội dung
            # Thường shape.name chứa "Title" hoặc là shape đầu tiên của slide
            is_title = "Title" in shape.name or shape == slide.shapes.title
            
            target_font = FONT_SANS if is_title else FONT_ROMAN
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = target_font
                    # Đảm bảo font được áp dụng (đôi khi pptx cần set ascii và hAnsi)
                    # Tuy nhiên bản 1.0+ xử lý khá tốt
                    
                    # Một số tùy chỉnh đặc biệt cho slide đầu (index 0)
                    if slide_idx == 0:
                        if is_title:
                            run.font.size = Pt(18)
                            run.font.bold = True
                        else:
                            run.font.size = Pt(10)
    
    prs.save(output_path)
    print(f"Done! Saved to {output_path}")

if __name__ == "__main__":
    input_pptx = "Slide_HUST.pptx"
    output_pptx = "Slide_HUST_Fixed.pptx"
    try:
        fix_fonts(input_pptx, output_pptx)
    except Exception as e:
        print(f"Error: {e}")
