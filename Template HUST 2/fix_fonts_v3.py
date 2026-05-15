from pptx import Presentation
from pptx.util import Pt

# Font chuẩn Mac
FONT_SANS = "LM Sans 10"

def fix_fonts_v3(input_path, output_path):
    prs = Presentation(input_path)
    slide_height = prs.slide_height
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Applying exact HUST specs to Slide {slide_idx + 1}...")
        
        for shape in slide.shapes:
            if shape.shape_type == 6: # Group
                shapes_to_process = shape.shapes
            else:
                shapes_to_process = [shape]
                
            for s in shapes_to_process:
                if not s.has_text_frame:
                    continue
                
                # Logic phân loại dựa trên vị trí và kích thước slide
                rel_top = s.top / slide_height
                
                for paragraph in s.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONT_SANS
                        
                        if slide_idx == 0:
                            # Title Slide specs từ beamerthemeHUST.sty
                            if rel_top < 0.5: # Tiêu đề chính
                                run.font.size = Pt(18)
                                run.font.bold = True
                            else: # Info SV/GV
                                run.font.size = Pt(10)
                        else:
                            # Content Slides
                            if rel_top < 0.15: # Frame Title
                                run.font.size = Pt(16)
                                run.font.bold = True
                            else: # Body text (size11.clo)
                                run.font.size = Pt(11)
                                
    prs.save(output_path)
    print(f"Đã xong bản v3 cực chuẩn tại {output_path}")

if __name__ == "__main__":
    fix_fonts_v3("Slide_HUST.pptx", "Slide_HUST_Fixed_v3.pptx")
