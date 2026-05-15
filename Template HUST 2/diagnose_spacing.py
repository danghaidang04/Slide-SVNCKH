"""
Kiểm tra sự khác biệt về kích thước giữa Arial và LM Sans 10
bằng cách so sánh text box width vs chiều rộng cần thiết.
"""
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

prs = Presentation("Slide_HUST.pptx")

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}

# Kiểm tra xem có run nào đã có character spacing không
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # Kiểm tra kích thước text box
        w_cm = shape.width / 360000 if shape.width else 0
        
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                
                # Kiểm tra character spacing hiện tại
                rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                spc = rPr.get('spc') if rPr is not None else None
                
                if spc:
                    text = run.text[:30]
                    print(f"  Slide {slide_idx+1} | spc={spc} | box_w={w_cm:.1f}cm | '{text}'")

print("\n--- Checking text box sizes ---")
for slide_idx, slide in enumerate(prs.slides):
    print(f"\nSlide {slide_idx+1}:")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        total_text = ""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                total_text += run.text
        
        if not total_text.strip():
            continue
        
        w_cm = shape.width / 360000 if shape.width else 0
        h_cm = shape.height / 360000 if shape.height else 0
        chars = len(total_text)
        
        # Tính mật độ ký tự / chiều rộng
        if w_cm > 0:
            density = chars / w_cm
            if density > 15:  # Nghi ngờ bị chật
                print(f"  ⚠️  '{shape.name}' w={w_cm:.1f}cm chars={chars} density={density:.0f}c/cm text='{total_text[:40]}'")
