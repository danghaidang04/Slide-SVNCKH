"""Inspect the original PPTX - safe version."""
from pptx import Presentation

prs = Presentation("Slide_HUST.pptx")

slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"Slide dimensions: {slide_w/914400*2.54:.2f} x {slide_h/914400*2.54:.2f} cm")
print(f"Beamer 16:9 dimensions: 16.00 x 9.00 cm")
print(f"Scale ratio (PPTX/Beamer): {(slide_w/914400*2.54)/16.0:.4f}")
print("="*80)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*80}")
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        top_cm = shape.top / 360000 if shape.top else 0
        left_cm = shape.left / 360000 if shape.left else 0
        
        print(f"\n  Shape: '{shape.name}' | Pos: ({left_cm:.1f}, {top_cm:.1f}) cm")
        
        seen_sizes = {}
        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                text_preview = run.text[:80].replace('\n', '\\n')
                font_name = run.font.name or "(inherited)"
                font_size = f"{run.font.size.pt:.1f}pt" if run.font.size else "(inherited)"
                font_bold = run.font.bold
                
                try:
                    font_color = str(run.font.color.rgb)
                except:
                    font_color = "(inherited)"
                
                key = f"{font_name}|{font_size}|{font_bold}"
                if key not in seen_sizes:
                    seen_sizes[key] = []
                seen_sizes[key].append(text_preview)
        
        for key, texts in seen_sizes.items():
            parts = key.split("|")
            sample = texts[0] if len(texts) == 1 else f"{texts[0]}... ({len(texts)} runs)"
            print(f"    Font: {parts[0]} | Size: {parts[1]} | Bold: {parts[2]}")
            print(f"    Sample: \"{sample}\"")
